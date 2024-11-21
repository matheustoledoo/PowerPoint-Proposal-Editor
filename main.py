import os
import sys
from PyQt5 import QtWidgets, QtCore
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import comtypes.client  # Para exportar para PDF (funciona apenas no Windows)

# Função para aplicar formatação
def aplicar_formatacao(paragraph, fonte="Codec Pro", tamanho=24, cor=(0, 0, 0)):
    for run in paragraph.runs:
        run.font.name = fonte
        run.font.size = Pt(tamanho)
        run.font.color.rgb = RGBColor(*cor)

# Função para exportar para PDF
def salvar_como_pdf(pptx_path, pdf_path):
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    try:
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        presentation.SaveAs(pdf_path, 32)  # 32 é o formato para PDF
        presentation.Close()
        return True
    except Exception as e:
        print(f"Erro ao salvar como PDF: {e}")
        return False
    finally:
        powerpoint.Quit()

# Funções restantes (manipulação do PowerPoint)...

# Classe da interface gráfica
class PowerPointEditor(QtWidgets.QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Editor de Propostas PowerPoint")
        self.resize(1000, 750)
        self.initUI()

    def initUI(self):
        layout = QtWidgets.QVBoxLayout(self)

        self.setStyleSheet("""
            QWidget { background-color: #f9f9f9; font-family: Arial, sans-serif; font-size: 14px; }
            QLabel { font-size: 16px; color: #333; font-weight: bold; }
            QListWidget { background-color: #fff; border: 1px solid #ccc; border-radius: 5px; padding: 5px; }
            QLineEdit, QTextEdit { border: 1px solid #ccc; padding: 5px; border-radius: 4px; font-size: 14px; }
            QPushButton { background-color: #0078D7; color: #fff; font-size: 14px; padding: 8px; border-radius: 4px; font-weight: bold; }
        """)

        titulo = QtWidgets.QLabel("Editor de Propostas PowerPoint")
        titulo.setAlignment(QtCore.Qt.AlignCenter)
        titulo.setStyleSheet("font-size: 22px; font-weight: bold; margin-bottom: 20px;")
        layout.addWidget(titulo)

        layout.addWidget(QtWidgets.QLabel("Selecione o arquivo PowerPoint:"))
        self.files_list = QtWidgets.QListWidget()
        self.load_files()
        layout.addWidget(self.files_list)

        tabs = QtWidgets.QTabWidget()
        layout.addWidget(tabs)

        tab1 = QtWidgets.QWidget()
        tab1_layout = QtWidgets.QFormLayout()
        self.nome_cliente = QtWidgets.QLineEdit(placeholderText="Digite o nome do cliente")
        self.valor_servico = QtWidgets.QLineEdit(placeholderText="Ex.: 10000")
        self.valor_mobilizacao = QtWidgets.QLineEdit(placeholderText="Ex.: 2000")
        tab1_layout.addRow("Nome do Cliente:", self.nome_cliente)
        tab1_layout.addRow("Valor do Serviço:", self.valor_servico)
        tab1_layout.addRow("Valor da Mobilização:", self.valor_mobilizacao)
        tab1.setLayout(tab1_layout)
        tabs.addTab(tab1, "Informações Gerais")

        tab2 = QtWidgets.QWidget()
        tab2_layout = QtWidgets.QFormLayout()
        self.objetos = QtWidgets.QTextEdit(placeholderText="Digite os objetos (um por linha)")
        self.escopo = QtWidgets.QTextEdit(placeholderText="Digite os itens de escopo (um por linha)")
        self.campo = QtWidgets.QTextEdit(placeholderText="Digite quem irá ao campo (um por linha)")
        self.processamento = QtWidgets.QTextEdit(placeholderText="Digite quem fará o processamento (um por linha)")
        self.equipamentos = QtWidgets.QTextEdit(placeholderText="Digite os equipamentos envolvidos (um por linha)")
        tab2_layout.addRow("Objetos:", self.objetos)
        tab2_layout.addRow("Escopo:", self.escopo)
        tab2_layout.addRow("Campo:", self.campo)
        tab2_layout.addRow("Processamento:", self.processamento)
        tab2_layout.addRow("Equipamentos:", self.equipamentos)
        tab2.setLayout(tab2_layout)
        tabs.addTab(tab2, "Listas")

        tab3 = QtWidgets.QWidget()
        tab3_layout = QtWidgets.QVBoxLayout()
        self.texto_slide11 = QtWidgets.QTextEdit(placeholderText="Texto para Slide 11")
        tab3_layout.addWidget(self.texto_slide11)
        tab3.setLayout(tab3_layout)
        tabs.addTab(tab3, "Texto Slide 11")

        form_layout = QtWidgets.QFormLayout()
        self.nome_arquivo = QtWidgets.QLineEdit(placeholderText="Ex.: Proposta")
        self.complemento_arquivo = QtWidgets.QLineEdit(placeholderText="Ex.: Revisado")
        form_layout.addRow("Nome do Arquivo:", self.nome_arquivo)
        form_layout.addRow("Complemento do Arquivo:", self.complemento_arquivo)
        layout.addLayout(form_layout)

        self.process_button = QtWidgets.QPushButton("Gerar Arquivo")
        self.process_button.clicked.connect(self.processar_arquivo)
        layout.addWidget(self.process_button)

    def load_files(self):
        if not os.path.exists("files"):
            os.makedirs("files")
        files = [f for f in os.listdir("files") if f.endswith(".pptx")]
        self.files_list.addItems(files)

    def processar_arquivo(self):
        try:
            arquivo = self.files_list.currentItem().text()
        except AttributeError:
            QtWidgets.QMessageBox.warning(self, "Erro", "Nenhum arquivo selecionado.")
            return

        caminho_arquivo = os.path.join("files", arquivo)
        prs = Presentation(caminho_arquivo)

        nome_cliente = self.nome_cliente.text()
        valor_servico = self.valor_servico.text()
        valor_mobilizacao = self.valor_mobilizacao.text()
        objetos = self.objetos.toPlainText().splitlines()
        escopo = self.escopo.toPlainText().splitlines()
        campo = self.campo.toPlainText().splitlines()
        processamento = self.processamento.toPlainText().splitlines()
        equipamentos = self.equipamentos.toPlainText().splitlines()
        texto_slide11 = self.texto_slide11.toPlainText()

        substituir_valores_marcadores(prs.slides[1], "{", nome_cliente)
        adicionar_objetos_dinamicos(prs.slides[2], objetos)
        adicionar_escopo_dinamicos(prs.slides[3], escopo)
        adicionar_lista_incremental(prs.slides[7], "Campo", campo)
        adicionar_lista_incremental(prs.slides[7], "Processamento", processamento)
        adicionar_lista_equipamentos(prs.slides[8], ":", equipamentos)
        substituir_valores_marcadores(prs.slides[10], "{", valor_servico)
        substituir_valores_marcadores(prs.slides[10], "}", valor_mobilizacao)
        substituir_texto_slide(prs.slides[11], texto_slide11)

        caminho_salvar, _ = QtWidgets.QFileDialog.getSaveFileName(self, "Salvar Arquivo", f"{self.nome_arquivo.text()}_{self.complemento_arquivo.text()}", "PowerPoint Files (*.pptx);;PDF Files (*.pdf)")
        if caminho_salvar:
            if caminho_salvar.endswith(".pdf"):
                prs.save("temp.pptx")
                if salvar_como_pdf("temp.pptx", caminho_salvar):
                    QtWidgets.QMessageBox.information(self, "Sucesso", "Arquivo salvo como PDF com sucesso!")
                else:
                    QtWidgets.QMessageBox.warning(self, "Erro", "Não foi possível salvar como PDF.")
                os.remove("temp.pptx")
            else:
                prs.save(caminho_salvar)
                QtWidgets.QMessageBox.information(self, "Sucesso", "Arquivo salvo com sucesso!")

# Executar aplicação
if __name__ == "__main__":
    app = QtWidgets.QApplication(sys.argv)
    editor = PowerPointEditor()
    editor.show()
    sys.exit(app.exec_())
